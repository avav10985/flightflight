from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

FONT = "Microsoft JhengHei"
TITLE_COLOR = RGBColor(0x1F, 0x3A, 0x5F)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x66, 0x66, 0x66)


def set_font(run, size=18, bold=False, color=TEXT_COLOR):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Ensure East Asian font applied
    rPr = run._r.get_or_add_rPr()
    from lxml import etree
    ea = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", FONT)


def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_font(r, size=32, bold=True, color=TITLE_COLOR)
    # Underline bar
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = TITLE_COLOR
    line.line.fill.background()


def add_bullets(slide, items, left=0.6, top=1.5, width=12.1, height=5.5, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, text, bold = item
        else:
            level, text, bold = 0, item, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        r = p.add_run()
        prefix = "  " * level + ("• " if level == 0 else "– ")
        r.text = prefix + text
        set_font(r, size=size - level * 2, bold=bold)
        p.space_after = Pt(6)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------- Slide 1: 封面 ----------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFA)
bg.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "DIY 四軸無人機"
set_font(r, size=54, bold=True, color=TITLE_COLOR)

tb = s.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "元件清單 + rx 電路原理"
set_font(r, size=32, color=GRAY)

tb = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "目標：GPS 路徑 + 超音波避障 + WiFi 影像 + 電腦端 AI"
set_font(r, size=20, color=TEXT_COLOR)


# ---------- Slide 2: 系統架構 ----------
s = prs.slides.add_slide(blank)
add_title(s, "系統架構")

# Drone block
def box(x, y, w, h, title, lines, fill=RGBColor(0xE8, 0xF0, 0xF8), border=TITLE_COLOR):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, size=18, bold=True, color=TITLE_COLOR)
    for ln in lines:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        set_font(r, size=14)

box(0.5, 1.5, 4.2, 5.3, "無人機端", [
    "ESP32 飛控",
    "MPU6050（6 軸姿態）",
    "BMP280（氣壓定高）",
    "NEO-M8N GPS",
    "HC-SR04 超音波 × 4（避障）",
    "NRF24L01+PA/LNA（接收）",
    "ESC × 4 → 無刷馬達 × 4",
    "",
    "ESP32-CAM（獨立 WiFi 影像）",
])

box(5.2, 1.5, 3.3, 5.3, "地面端", [
    "Arduino Nano",
    "NRF24L01+PA/LNA",
    "雙搖桿（手動備援）",
    "OLED 顯示遙測",
    "USB 連電腦",
], fill=RGBColor(0xFD, 0xF2, 0xE8))

box(9.0, 1.5, 3.8, 5.3, "電腦端", [
    "Python 接收指令",
    "YOLO 物體辨識",
    "路徑規劃",
    "發送 GOTO / MOVE 指令",
    "",
    "WiFi 直連 ESP32-CAM",
    "看即時影像",
], fill=RGBColor(0xE8, 0xF4, 0xE8))


# ---------- Slide 3: 你已有的元件 ----------
s = prs.slides.add_slide(blank)
add_title(s, "你已有的元件（不用買）")
add_bullets(s, [
    (0, "無刷馬達 × 4", True),
    (0, "ESC 電調 × 4", True),
    (0, "NRF24L01+PA/LNA 模組（多片）", True),
    (0, "Arduino Nano（地面端用）", True),
    (0, "機架（?）、電池（?）、螺旋槳（?）", True),
    (0, "", False),
    (0, "確認一下：機架尺寸、電池規格（幾 S 多少 mAh）、螺旋槳尺寸", False),
    (0, "這會影響接下來要不要買電源分配板 PDB 和 BEC", False),
], size=20)


# ---------- Slide 4: 必買 - 飛控核心 ----------
s = prs.slides.add_slide(blank)
add_title(s, "需要購買 ①：飛控核心（約 NT$450）")

rows = [
    ("元件", "型號", "用途", "價格"),
    ("主控板", "ESP32 DevKit V1", "跑 PID、感測器、NRF24、GPS", "250"),
    ("6 軸 IMU", "MPU6050 模組", "陀螺儀 + 加速度計，自穩必備", "60"),
    ("氣壓計", "BMP280 模組", "定高模式使用", "60"),
    ("蜂鳴器", "5V 有源蜂鳴器", "低電量 / 失聯警報", "20"),
    ("小計", "", "", "390"),
]
left = 0.8
top = 1.6
widths = [2.0, 3.5, 5.5, 1.5]
rowh = 0.55

cx = left
for i, (w, h) in enumerate(zip(widths, [rowh]*len(rows))):
    pass

y = top
for ri, row in enumerate(rows):
    x = left
    for ci, text in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(widths[ci]), Inches(rowh))
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TITLE_COLOR
            cell.line.color.rgb = TITLE_COLOR
        elif ri == len(rows) - 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
            cell.line.color.rgb = GRAY
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.line.color.rgb = GRAY
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        if ri == 0:
            set_font(r, size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        elif ri == len(rows) - 1:
            set_font(r, size=14, bold=True)
        else:
            set_font(r, size=14)
        x += widths[ci]
    y += rowh


# ---------- Slide 5: 必買 - 定位與避障 ----------
s = prs.slides.add_slide(blank)
add_title(s, "需要購買 ②：GPS + 避障（約 NT$560）")
add_bullets(s, [
    (0, "GPS 模組：NEO-M8N（NT$400）", True),
    (1, "支援 GPS + GLONASS，定位精度 2.5m", False),
    (1, "便宜版 NEO-6M 只要 NT$150，但冷啟動慢、精度差", False),
    (0, "超音波感測器：HC-SR04 × 4（NT$40 × 4 = NT$160）", True),
    (1, "前、後、左、右各一顆", False),
    (1, "射程 2~4m、便宜好上手", False),
    (1, "注意：戶外風噪會影響，可先當「警示」不做硬避障", False),
    (0, "若日後要升級，可換 VL53L1X 雷射 ToF（4 顆 NT$600）", False),
], size=18)


# ---------- Slide 6: 必買 - 影像 ----------
s = prs.slides.add_slide(blank)
add_title(s, "需要購買 ③：影像模組（約 NT$300）")
add_bullets(s, [
    (0, "ESP32-CAM（AI-Thinker 版）：NT$250", True),
    (1, "內建 OV2640 相機 + WiFi", False),
    (1, "電腦用瀏覽器或 Python 直連就能看畫面", False),
    (0, "外接 WiFi 天線（IPEX）：NT$50", True),
    (1, "原廠 PCB 天線只有 100m 左右，外接可延伸到 300~500m", False),
    (0, "FTDI USB 轉 TTL（燒錄用）：NT$80", True),
    (1, "ESP32-CAM 沒有 USB，燒錄必備", False),
    (1, "飛控的 ESP32 DevKit 有 USB 所以不用", False),
    (0, "", False),
    (0, "注意：ESP32-CAM 獨立運作，不跟飛控通訊，不會互相影響", False),
], size=18)


# ---------- Slide 7: 必買 - 其他 ----------
s = prs.slides.add_slide(blank)
add_title(s, "需要購買 ④：其他配件（約 NT$400）")
add_bullets(s, [
    (0, "電源分配板 PDB（含 5V BEC）：NT$150", True),
    (1, "給飛控和感測器穩定 5V 供電", False),
    (0, "XT60 接頭 + 電源線：NT$80", True),
    (0, "杜邦線、排針、束帶、熱縮套：NT$150", True),
    (0, "OLED 0.96 吋顯示器（地面端用）：NT$80", True),
    (1, "顯示飛機狀態、GPS 座標、電量", False),
    (0, "", False),
    (0, "搖桿模組（地面端，你的 TX 已經有）", False),
    (0, "Arduino Nano（地面端，你已有）", False),
], size=18)


# ---------- Slide 8: BOM 總表 ----------
s = prs.slides.add_slide(blank)
add_title(s, "BOM 總表")

rows = [
    ("分類", "內容", "小計 (NT$)"),
    ("飛控核心", "ESP32 + MPU6050 + BMP280 + 蜂鳴器", "390"),
    ("定位避障", "NEO-M8N + HC-SR04 × 4", "560"),
    ("影像模組", "ESP32-CAM + 天線 + FTDI", "380"),
    ("配件", "PDB + XT60 + 杜邦線 + OLED", "460"),
    ("總計", "（馬達/ESC/電池/機架/螺旋槳已有）", "1,790"),
]
widths = [2.5, 7.5, 2.5]
rowh = 0.6
left = 0.4; top = 1.6
y = top
for ri, row in enumerate(rows):
    x = left
    for ci, text in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(widths[ci]), Inches(rowh))
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TITLE_COLOR
        elif ri == len(rows) - 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xCC)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.line.color.rgb = GRAY
        tf = cell.text_frame
        tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        if ri == 0:
            set_font(r, size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        elif ri == len(rows) - 1:
            set_font(r, size=16, bold=True, color=ACCENT)
        else:
            set_font(r, size=15)
        x += widths[ci]
    y += rowh

tb = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0))
p = tb.text_frame.paragraphs[0]
r = p.add_run()
r.text = "✓ 預算目標 NT$2000~3000 內完成採購（馬達/電池自備的前提下）"
set_font(r, size=18, bold=True, color=ACCENT)


# ---------- Slide 9: rx 電路整體功能 ----------
s = prs.slides.add_slide(blank)
add_title(s, "rx.JPG 電路原理：這是什麼？")
add_bullets(s, [
    (0, "這張圖是「有刷馬達」四軸的馬達驅動電路", True),
    (1, "4 顆小型有刷直流馬達（像玩具小四軸那種）", False),
    (1, "不適合你的無刷馬達 → 無刷馬達要用 ESC", False),
    (0, "", False),
    (0, "每一顆馬達的驅動電路包含 4 個元件：", True),
    (1, "1. N 通道 MOSFET（IRFZ44N）→ 當作電子開關", False),
    (1, "2. 100Ω 閘極電阻 → 限制瞬間電流", False),
    (1, "3. 10kΩ 下拉電阻 → 安全接地", False),
    (1, "4. 1N4007 飛輪二極體 → 保護電路免於反電動勢", False),
    (0, "", False),
    (0, "下面幾頁逐個說明每個元件的電子學原理", False),
], size=18)


# ---------- Slide 10: MOSFET 低壓端開關 ----------
s = prs.slides.add_slide(blank)
add_title(s, "① MOSFET：電子開關")
add_bullets(s, [
    (0, "IRFZ44N 是 N 通道 MOSFET（功率電晶體）", True),
    (1, "三隻腳：Gate (G)、Drain (D)、Source (S)", False),
    (1, "最大可過 49A，阻抗低發熱小，適合切換馬達", False),
    (0, "接法：「低壓端開關」（Low-Side Switch）", True),
    (1, "馬達正極 → 接電池 +V", False),
    (1, "馬達負極 → 接 MOSFET 的 Drain", False),
    (1, "MOSFET Source → 接地（GND）", False),
    (1, "Arduino 用 Gate 控制開關", False),
    (0, "運作方式（開關狀態）：", True),
    (1, "Gate 電壓 > 4V（HIGH）：D-S 導通 → 馬達通電轉動", False),
    (1, "Gate 電壓 ≈ 0V（LOW）：D-S 截止 → 馬達斷電停止", False),
], size=17)


# ---------- Slide 11: PWM 調速 ----------
s = prs.slides.add_slide(blank)
add_title(s, "② 為什麼能控制「轉速」？PWM 原理")
add_bullets(s, [
    (0, "Arduino analogWrite() 輸出 PWM 方波（預設約 490Hz）", True),
    (0, "Duty Cycle（佔空比）= HIGH 時間 / 週期", True),
    (1, "0% Duty：MOSFET 一直關 → 馬達不轉", False),
    (1, "50% Duty：一半時間通電 → 馬達半速", False),
    (1, "100% Duty：一直通 → 馬達全速", False),
    (0, "原理：因為切換速度很快（每秒 490 次）", True),
    (1, "馬達的機械慣性來不及跟著開關", False),
    (1, "馬達「感覺到」的是平均電壓", False),
    (1, "平均電壓 = 電池電壓 × Duty Cycle", False),
    (0, "", False),
    (0, "所以：改變 PWM 值 → 改變平均電壓 → 改變轉速", True),
], size=18)


# ---------- Slide 12: 閘極電阻 ----------
s = prs.slides.add_slide(blank)
add_title(s, "③ 100Ω 閘極電阻的作用")
add_bullets(s, [
    (0, "MOSFET 的 Gate 其實是一個電容（約 1~10 nF）", True),
    (1, "要導通 = 把這個電容充飽到 > 4V", False),
    (1, "要關閉 = 把這個電容放電到 ≈ 0V", False),
    (0, "如果沒有閘極電阻：", True),
    (1, "Arduino 瞬間要灌 / 吸很大電流（可能 >100mA）", False),
    (1, "會超過 Arduino 腳位的 40mA 上限 → 燒 MCU", False),
    (1, "開關邊緣會產生振鈴（ringing）電磁干擾", False),
    (0, "加上 100Ω 之後：", True),
    (1, "限制瞬間電流 = 5V / 100Ω = 50mA（安全）", False),
    (1, "稍微減慢開關速度，減少振鈴", False),
    (1, "代價：切換不那麼陡峭，MOSFET 在切換時會發熱一點點", False),
], size=17)


# ---------- Slide 13: 下拉電阻 ----------
s = prs.slides.add_slide(blank)
add_title(s, "④ 10kΩ 下拉電阻的作用（安全！）")
add_bullets(s, [
    (0, "問題情境：Arduino 重開機 / 程式沒跑 / 線鬆掉", True),
    (1, "此時 Gate 是「浮接」狀態（floating）", False),
    (1, "浮接腳位會受環境雜訊感應出不定電壓", False),
    (1, "可能意外到達 4V → MOSFET 誤導通 → 馬達亂轉！", False),
    (0, "解法：在 Gate 和 GND 之間加 10kΩ 電阻", True),
    (1, "Arduino 沒在輸出時，Gate 被拉到 0V", False),
    (1, "MOSFET 保持關閉 → 馬達安全停止", False),
    (0, "為什麼選 10kΩ？", True),
    (1, "太小（如 1kΩ）→ 正常 HIGH 時浪費電、會分壓降低 Gate 電壓", False),
    (1, "太大（如 1MΩ）→ 放電太慢，Gate 電壓降不下來", False),
    (1, "10kΩ 是經驗值的甜蜜點", False),
], size=17)


# ---------- Slide 14: 飛輪二極體 ----------
s = prs.slides.add_slide(blank)
add_title(s, "⑤ 1N4007 飛輪二極體（Flyback Diode）")
add_bullets(s, [
    (0, "馬達是電感（線圈）—— 記住這點最重要！", True),
    (0, "電感的特性：不允許電流瞬間改變", True),
    (1, "MOSFET 關閉的瞬間，電流突然被切斷", False),
    (1, "電感會自己產生「反電動勢」，電壓可能飆到 100V 以上", False),
    (1, "這個高壓會擊穿 MOSFET → 燒機", False),
    (0, "解法：在馬達兩端反向並聯一顆二極體", True),
    (1, "陰極（有線那一端）→ 接 +V", False),
    (1, "陽極 → 接 MOSFET Drain（馬達負極）", False),
    (1, "正常工作時二極體反向，不影響電路", False),
    (1, "MOSFET 關閉瞬間，電感感應電壓讓二極體正向導通", False),
    (1, "電感能量透過二極體安全釋放 → MOSFET 得救", False),
    (0, "1N4007 規格 1A / 1000V，便宜又耐用", False),
], size=15)


# ---------- Slide 15: 為何不能用於無刷 ----------
s = prs.slides.add_slide(blank)
add_title(s, "為什麼這個電路不能控制「無刷馬達」？")
add_bullets(s, [
    (0, "有刷 vs 無刷：根本不同的東西", True),
    (1, "有刷馬達：兩條線，直接給電就轉（像玩具車）", False),
    (1, "無刷馬達：三條線，需要按順序切換三相電流", False),
    (0, "無刷馬達需要：", True),
    (1, "6 顆 MOSFET 組成三相橋", False),
    (1, "專用晶片計算換相時序（或感測轉子位置）", False),
    (1, "換相頻率從每秒幾十到幾萬次", False),
    (0, "解法：買現成的「電子調速器 ESC」", True),
    (1, "ESC 裡面已經包含三相橋 + 換相邏輯", False),
    (1, "對飛控來說，只需要給它一條 PWM 訊號（像控制伺服那樣）", False),
    (1, "1000µs = 停止，2000µs = 全速", False),
    (0, "", False),
    (0, "所以你已經有 ESC → 不需要自己焊 MOSFET 電路", True),
], size=17)


# ---------- Slide 16: 結語 ----------
s = prs.slides.add_slide(blank)
add_title(s, "下一步")
add_bullets(s, [
    (0, "採購清單（總預算 ~NT$1800）", True),
    (1, "ESP32 + MPU6050 + BMP280 + 蜂鳴器", False),
    (1, "NEO-M8N GPS + HC-SR04 × 4", False),
    (1, "ESP32-CAM + 外接天線 + FTDI", False),
    (1, "PDB + XT60 + 杜邦線 + OLED", False),
    (0, "到貨後依序進行：", True),
    (1, "Step 1：ESP32 + MPU6050 單獨測試（讀姿態角）", False),
    (1, "Step 2：加入 NRF24 接收，手動搖桿控制四顆 ESC", False),
    (1, "Step 3：綁上機架試飛（無 GPS、無避障）", False),
    (1, "Step 4：加入 GPS 定點懸停", False),
    (1, "Step 5：加入超音波避障邏輯", False),
    (1, "Step 6：加入 ESP32-CAM，測試電腦接收影像", False),
    (1, "Step 7：電腦端寫 AI，閉環控制", False),
], size=17)


out = r"d:\無人機\flightflight\drone_components.pptx"
prs.save(out)
print(f"Saved: {out}")
